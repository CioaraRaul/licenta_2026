# -*- coding: utf-8 -*-
"""Generează prezentarea PowerPoint pentru proiectul de licență AutoVault."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(ROOT, "screenshots")
OUT = os.path.join(ROOT, "Licenta-AutoMarketplace.pptx")

# Format 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Paletă de culori
NAVY = RGBColor(0x0F, 0x23, 0x4E)
NAVY_DARK = RGBColor(0x07, 0x14, 0x2F)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)


def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=20, color=TEXT, bullet_color=ACCENT,
                line_spacing=1.15, space_after=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None, page=None, total=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.08), fill=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.6),
             title, size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.82), Inches(10.5), Inches(0.4),
                 subtitle, size=16, color=RGBColor(0xCB, 0xD5, 0xE1))
    if page is not None and total is not None:
        add_text(slide, Inches(11.4), Inches(0.45), Inches(1.7), Inches(0.5),
                 f"{page} / {total}", size=16, bold=True,
                 color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.RIGHT)


def footer(slide):
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), fill=NAVY_DARK)
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             "Licență 2026  •  AutoVault Marketplace",
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(slide, Inches(9.5), Inches(7.15), Inches(3.3), Inches(0.3),
             "NestJS  •  React Router 7  •  TypeORM",
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.RIGHT)


def add_image_safe(slide, image_name, x, y, w, h, *, caption=None):
    path = os.path.join(SHOTS, image_name) if image_name else None
    border = add_rect(slide, x, y, w, h, fill=LIGHT, line=RGBColor(0xD1, 0xD5, 0xDB))
    if path and os.path.isfile(path):
        try:
            with Image.open(path) as im:
                iw, ih = im.size
            target_ratio = w / h
            img_ratio = iw / ih
            if img_ratio > target_ratio:
                new_w = w
                new_h = int(w / img_ratio)
                nx = x
                ny = y + (h - new_h) // 2
            else:
                new_h = h
                new_w = int(h * img_ratio)
                ny = y
                nx = x + (w - new_w) // 2
            slide.shapes.add_picture(path, nx, ny, width=new_w, height=new_h)
        except Exception as e:
            add_text(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.5),
                     f"[Eroare imagine: {e}]", size=12, color=MUTED)
    else:
        add_text(slide, x, y + h / 2 - Inches(0.4), w, Inches(0.4),
                 f"[Captură ecran: {image_name or 'placeholder'}]",
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + h / 2 + Inches(0.1), w, Inches(0.4),
                 "Inserați aici imaginea din aplicație",
                 size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, x, y + h + Inches(0.05), w, Inches(0.3),
                 caption, size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ---------- constructori slide-uri ----------

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_title(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_text(s, 0, Inches(0.7), SLIDE_W, Inches(0.5),
             "LUCRARE DE DIPLOMĂ  •  SESIUNEA IULIE 2026", size=20, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(1.2), SLIDE_W, Inches(0.4),
             "Universitatea din Oradea  •  Specializarea Calculatoare",
             size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    add_rect(s, Inches(5.92), Inches(2.05), Inches(1.5), Inches(0.1), fill=ACCENT)

    add_text(s, 0, Inches(2.4), SLIDE_W, Inches(1.0),
             "Platformă web pentru vânzarea și",
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(3.3), SLIDE_W, Inches(1.0),
             "cumpărarea automobilelor",
             size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.45), SLIDE_W, Inches(0.7),
             "utilizând NestJS și React",
             size=30, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)

    add_rect(s, 0, Inches(6.0), SLIDE_W, Inches(1.5), fill=NAVY_DARK)

    add_text(s, Inches(0.8), Inches(6.18), Inches(6), Inches(0.45),
             "ABSOLVENT", size=13, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(6.55), Inches(6), Inches(0.6),
             "Cioara Raul", size=24, bold=True, color=WHITE)

    add_text(s, Inches(7), Inches(6.18), Inches(5.8), Inches(0.45),
             "COORDONATOR ȘTIINȚIFIC", size=13, bold=True, color=ACCENT,
             align=PP_ALIGN.RIGHT)
    add_text(s, Inches(7), Inches(6.55), Inches(5.8), Inches(0.6),
             "Ș.L.DR.ING. PECHERLE GEORGE-DOMINIC", size=20, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT)


def slide_text(prs, title, bullets, *, page, total, subtitle=None, body_size=24):
    s = new_slide(prs)
    header(s, title, subtitle=subtitle, page=page, total=total)
    add_bullets(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(5.0),
                bullets, size=body_size, line_spacing=1.25, space_after=14)
    footer(s)
    return s


def slide_two_columns(prs, title, left_title, left_items, right_title, right_items,
                       *, page, total, subtitle=None, body_size=18):
    s = new_slide(prs)
    header(s, title, subtitle=subtitle, page=page, total=total)
    card_y = Inches(1.85)
    card_h = Inches(5.15)
    add_rect(s, Inches(0.54), card_y + Inches(0.05), Inches(6.1), card_h,
             fill=RGBColor(0xE5, 0xE7, 0xEB), line=None)
    add_rect(s, Inches(0.5), card_y, Inches(6.1), card_h,
             fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(s, Inches(0.5), card_y, Inches(6.1), Inches(0.8), fill=NAVY)
    add_text(s, Inches(0.5), card_y, Inches(6.1), Inches(0.8),
             left_title, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.85), card_y + Inches(1.0), Inches(5.5), Inches(4.0),
                left_items, size=body_size, line_spacing=1.2, space_after=10)
    add_rect(s, Inches(6.77), card_y + Inches(0.05), Inches(6.1), card_h,
             fill=RGBColor(0xE5, 0xE7, 0xEB), line=None)
    add_rect(s, Inches(6.73), card_y, Inches(6.1), card_h,
             fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
    add_rect(s, Inches(6.73), card_y, Inches(6.1), Inches(0.8), fill=NAVY)
    add_text(s, Inches(6.73), card_y, Inches(6.1), Inches(0.8),
             right_title, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(7.08), card_y + Inches(1.0), Inches(5.5), Inches(4.0),
                right_items, size=body_size, line_spacing=1.2, space_after=10)
    footer(s)


def slide_with_image(prs, title, image_name, bullets, *, page, total, subtitle=None,
                     image_left=False, image_caption=None):
    s = new_slide(prs)
    header(s, title, subtitle=subtitle, page=page, total=total)
    img_x = Inches(0.5) if image_left else Inches(6.85)
    bul_x = Inches(7.0) if image_left else Inches(0.7)
    add_image_safe(s, image_name, img_x, Inches(1.95), Inches(6.0), Inches(4.3),
                   caption=image_caption)
    add_bullets(s, bul_x, Inches(2.0), Inches(5.85), Inches(4.5),
                bullets, size=19, line_spacing=1.2, space_after=12)
    footer(s)


def slide_big_image(prs, title, image_name, *, page, total, subtitle=None, caption=None):
    s = new_slide(prs)
    header(s, title, subtitle=subtitle, page=page, total=total)
    add_image_safe(s, image_name, Inches(1.0), Inches(1.95), Inches(11.33), Inches(4.7),
                   caption=caption)
    footer(s)


def slide_tech_grid(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Stiva tehnologică completă",
           subtitle="Tehnologii moderne pentru un marketplace web complet",
           page=page, total=total)
    items = [
        ("PARTEA DE SERVER", "NestJS 11", "Cadru Node.js modular cu injecție de dependențe"),
        ("LIMBAJ", "TypeScript", "Tipare statice, decoratori și obiecte de transfer validate"),
        ("BAZĂ DE DATE", "SQLite + TypeORM", "Mapare relațională cu sincronizare automată"),
        ("AUTENTIFICARE", "JWT + OAuth2", "Jetoane de acces și reîmprospătare, Google și Facebook"),
        ("INTERFAȚA", "React 19 + RR 7", "Randare pe server cu React Router în mod cadru"),
        ("STARE", "Zustand 5", "Magazine persistente în memoria locală a browserului"),
        ("STILIZARE", "Tailwind CSS 4", "Utilitar-first cu jetoane de design"),
        ("MEDIA", "Cloudinary + Multer", "Încărcare optimizată de imagini"),
    ]
    cols = 4
    box_w = Inches(3.05)
    box_h = Inches(2.45)
    gap = Inches(0.12)
    x0 = Inches(0.42)
    y0 = Inches(1.85)
    for i, (badge, name, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + c * (box_w + gap)
        y = y0 + r * (box_h + Inches(0.2))
        add_rect(s, x + Inches(0.04), y + Inches(0.05), box_w, box_h,
                 fill=RGBColor(0xE5, 0xE7, 0xEB), line=None)
        add_rect(s, x, y, box_w, box_h, fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
        add_rect(s, x, y, box_w, Inches(0.55), fill=NAVY)
        add_text(s, x, y + Inches(0.1), box_w, Inches(0.4),
                 badge, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(0.78), box_w - Inches(0.3), Inches(0.7),
                 name, size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), y + Inches(1.6), box_w - Inches(0.5), Inches(0.9),
                 desc, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)


def slide_architecture(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Arhitectura generală a aplicației",
           subtitle="Fluxul datelor: client → interfață programatică → mapare relațională → bază de date",
           page=page, total=total)
    boxes = [
        ("CLIENT", "Browser (React 19, Vite, Tailwind CSS 4)",
         "Pagini randate pe server și client, magazine Zustand cu persistență în browser", BLUE),
        ("TRANSPORT", "Axios cu interceptori de rețea",
         "Jeton de acces automat, reîmprospătare automată la eroare 401, coadă de cereri", ACCENT),
        ("API", "Controlere NestJS cu gardieni de rute",
         "Conductă de validare globală, autentificare JWT, limitare cereri, CORS configurat", NAVY),
        ("DATE", "TypeORM → SQLite în mod WAL",
         "Entități și relații tipate, sincronizare automată a schemei în dezvoltare", GREEN),
    ]
    y = Inches(1.85)
    box_h = Inches(1.18)
    gap = Inches(0.1)
    for badge, title, desc, color in boxes:
        add_rect(s, Inches(0.74), y + Inches(0.04), Inches(11.9), box_h,
                 fill=RGBColor(0xE5, 0xE7, 0xEB), line=None)
        add_rect(s, Inches(0.7), y, Inches(11.9), box_h,
                 fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
        add_rect(s, Inches(0.7), y, Inches(2.1), box_h, fill=color)
        add_text(s, Inches(0.7), y, Inches(2.1), box_h,
                 badge, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.0), y + Inches(0.18), Inches(9.55), Inches(0.5),
                 title, size=22, bold=True, color=NAVY)
        add_text(s, Inches(3.0), y + Inches(0.7), Inches(9.55), Inches(0.45),
                 desc, size=15, color=MUTED)
        y += box_h + gap
    footer(s)


def slide_erd(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Modelul de date",
           subtitle="Cele 9 entități principale și câmpurile lor cheie",
           page=page, total=total)

    col_w = Inches(3.85)
    row_h = Inches(1.55)
    gap_x = Inches(0.38)
    gap_y = Inches(0.30)
    x0 = Inches(0.45)
    y0 = Inches(1.82)
    cols = [x0, x0 + col_w + gap_x, x0 + 2 * (col_w + gap_x)]
    rows = [y0, y0 + row_h + gap_y, y0 + 2 * (row_h + gap_y)]

    LINE_C = RGBColor(0x9C, 0xA3, 0xAF)

    def entity(x, y, w, h, name, fields, color=NAVY):
        add_rect(s, x + Inches(0.04), y + Inches(0.05), w, h,
                 fill=RGBColor(0xE5, 0xE7, 0xEB), line=None)
        add_rect(s, x, y, w, h, fill=WHITE, line=color)
        add_rect(s, x, y, w, Inches(0.52), fill=color)
        add_text(s, x + Inches(0.18), y + Inches(0.06), w - Inches(0.25), Inches(0.42),
                 name, size=19, bold=True, color=WHITE)
        for i, f in enumerate(fields):
            add_text(s, x + Inches(0.18), y + Inches(0.62) + i * Inches(0.3),
                     w - Inches(0.25), Inches(0.3),
                     f, size=13, color=TEXT)

    def hline(gx, row_y, label):
        cy = row_y + row_h / 2
        add_rect(s, gx, cy - Inches(0.02), gap_x, Inches(0.04), fill=LINE_C, line=None)
        add_text(s, gx, cy - Inches(0.24), gap_x, Inches(0.22),
                 label, size=9, bold=True, color=LINE_C, align=PP_ALIGN.CENTER)

    def vline(col_x, gy, label):
        cx = col_x + col_w / 2
        add_rect(s, cx - Inches(0.02), gy, Inches(0.04), gap_y, fill=LINE_C, line=None)
        add_text(s, cx + Inches(0.06), gy, Inches(0.32), gap_y,
                 label, size=9, bold=True, color=LINE_C)

    # Rândul 1
    entity(cols[0], rows[0], col_w, row_h, "Cumpărător",
           ["soldDisponibil, monedă", "totalCheltuit, alerteVehicule"], color=BLUE)
    entity(cols[1], rows[0], col_w, row_h, "Utilizator",
           ["id, adresăEmail, parolă", "rol, prenume, numeFamilie",
            "emailVerificat, datăCreare"], color=NAVY)
    entity(cols[2], rows[0], col_w, row_h, "Vehicul",
           ["marcă, model, an", "preț, kilometraj, combustibil",
            "imagini, vizualizări, stare"], color=GREEN)

    # Rândul 2
    entity(cols[0], rows[1], col_w, row_h, "Vânzător",
           ["soldDisponibil, IBAN", "cotăComision, evaluare"], color=BLUE)
    entity(cols[1], rows[1], col_w, row_h, "Conversație",
           ["idCumpărător, idVânzător, idVehicul",
            "necititeCumpărător, necititeVânzător",
            "ultimMesaj, datăUltimMesaj"], color=ACCENT)
    entity(cols[2], rows[1], col_w, row_h, "Ofertă",
           ["sumă, stare, mesaj",
            "expirăLa, motivRespingere",
            "idCumpărător, idVehicul"], color=ACCENT)

    # Rândul 3
    entity(cols[0], rows[2], col_w, row_h, "VehiculSalvat",
           ["idUtilizator, idVehicul", "unic (idUtilizator, idVehicul)"], color=BLUE)
    entity(cols[1], rows[2], col_w, row_h, "Mesaj",
           ["conținut, stare", "idTrimițător, idConversație",
            "datăCreare"], color=ACCENT)
    entity(cols[2], rows[2], col_w, row_h, "Tranzacție",
           ["sumă, tip, stare", "idUtilizator, idReferință",
            "descriere, datăCreare"], color=GREEN)

    footer(s)


def slide_modules(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Module ale serverului NestJS",
           subtitle="9 module principale, fiecare cu controler, serviciu, entitate și obiect de transfer",
           page=page, total=total)
    modules = [
        ("utilizatori / autentificare", "Înregistrare, autentificare, JWT, OAuth, verificare email, resetare parolă"),
        ("vehicule", "Operații complete pe anunțuri cu 40+ câmpuri tehnice; filtrare și sortare avansată"),
        ("oferte", "Plasare oferte de cumpărători, acceptare/respingere de vânzători"),
        ("mesaje", "Conversații persistente cumpărător-vânzător-vehicul cu urmărire necitite"),
        ("vehicule-salvate", "Listă de favorite per cumpărător cu constrângere de unicitate"),
        ("portofel", "Carduri, depozite, retrageri, sume blocate pentru oferte"),
        ("tablou-de-bord", "Agregare metrici în funcție de rol (CUMPĂRĂTOR / VÂNZĂTOR)"),
        ("încărcare", "Încărcare imagini pe Cloudinary cu middleware Multer"),
        ("comun", "Gardieni de rute, decoratori, filtre de excepții, configurare limitare cereri"),
    ]
    y0 = Inches(1.75)
    for i, (name, desc) in enumerate(modules):
        r, c = divmod(i, 3)
        x = Inches(0.5) + c * Inches(4.2)
        y = y0 + r * Inches(1.65)
        add_rect(s, x, y, Inches(4.0), Inches(1.5),
                 fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
        add_text(s, x + Inches(0.2), y + Inches(0.18), Inches(3.8), Inches(0.4),
                 name, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), y + Inches(0.68), Inches(3.8), Inches(0.8),
                 desc, size=13, color=MUTED)
    footer(s)


def slide_four_roles(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Tipuri de utilizatori",
           subtitle="Patru roluri cu acces granular pe resurse",
           page=page, total=total)
    roles = [
        ("CUMPĂRĂTOR", BLUE, "rol pentru achiziții",
         ["Caută și filtrează vehicule.",
          "Salvează favorite și compară mașini.",
          "Plasează oferte pe vehicule.",
          "Comunică prin mesagerie.",
          "Portofel cu sume blocate."]),
        ("VÂNZĂTOR", GREEN, "rol pentru vânzări",
         ["Adaugă anunțuri detaliate.",
          "Editează și marchează ca VÂNDUT.",
          "Acceptă sau respinge oferte.",
          "Comunică cu cumpărătorii.",
          "Vede câștiguri și comisioane."]),
        ("VIZITATOR", ACCENT, "rol pentru navigare",
         ["Acces public la anunțuri.",
          "Vede detalii complete.",
          "Nu poate salva sau oferta.",
          "Invitat să se autentifice."]),
        ("ADMINISTRATOR", NAVY, "rol de administrare (în dezvoltare)",
         ["Rapoarte agregate.",
          "Pregătit la nivel de entitate.",
          "Implementare planificată."]),
    ]
    box_w = Inches(6.2)
    box_h = Inches(2.45)
    gap_x = Inches(0.13)
    gap_y = Inches(0.18)
    x0 = Inches(0.45)
    y0 = Inches(1.85)
    for i, (badge, color, sub, items) in enumerate(roles):
        r, c = divmod(i, 2)
        x = x0 + c * (box_w + gap_x)
        y = y0 + r * (box_h + gap_y)
        add_rect(s, x + Inches(0.04), y + Inches(0.05), box_w, box_h,
                 fill=RGBColor(0xE5, 0xE7, 0xEB), line=None)
        add_rect(s, x, y, box_w, box_h, fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
        add_rect(s, x, y, Inches(2.0), box_h, fill=color)
        add_text(s, x, y + Inches(0.45), Inches(2.0), Inches(0.55),
                 badge, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.1), Inches(2.0), Inches(0.4),
                 sub, size=12, color=RGBColor(0xE5, 0xE7, 0xEB),
                 align=PP_ALIGN.CENTER)
        add_bullets(s, x + Inches(2.2), y + Inches(0.25),
                    box_w - Inches(2.35), box_h - Inches(0.4),
                    items, size=16, line_spacing=1.22, space_after=8)
    footer(s)


def slide_auth(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Autentificare și înregistrare",
           subtitle="Autentificare locală și prin Google / Facebook, cu protecție împotriva atacurilor de forță brută",
           page=page, total=total)

    # Două imagini side-by-side: login (stânga) și register (dreapta)
    img_y = Inches(1.85)
    img_h = Inches(3.85)
    img_w = Inches(6.2)
    gap = Inches(0.2)
    x_left = Inches(0.4)
    x_right = x_left + img_w + gap

    add_image_safe(s, "01-login.png", x_left, img_y, img_w, img_h,
                   caption="Autentificare — email / parolă sau cont Google / Facebook")
    add_image_safe(s, "02-register.png", x_right, img_y, img_w, img_h,
                   caption="Înregistrare — creare cont nou cu validare câmpuri")

    # Bandă de caracteristici cheie la baza slide-ului
    features = [
        ("JWT", "Acces 15 min + reîmprospătare 7 zile"),
        ("scrypt", "Parole cifrate cu salt aleatoriu"),
        ("Throttling", "Max 5 încercări / minut"),
        ("Email", "Verificare adresă la înregistrare"),
    ]
    feat_w = Inches(3.05)
    feat_gap = Inches(0.14)
    x_feat = Inches(0.42)
    y_feat = Inches(6.15)
    for badge, desc in features:
        add_rect(s, x_feat, y_feat, feat_w, Inches(0.75),
                 fill=WHITE, line=RGBColor(0xE5, 0xE7, 0xEB))
        add_rect(s, x_feat, y_feat, Inches(1.1), Inches(0.75), fill=NAVY)
        add_text(s, x_feat, y_feat, Inches(1.1), Inches(0.75),
                 badge, size=14, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x_feat + Inches(1.18), y_feat + Inches(0.12),
                 feat_w - Inches(1.25), Inches(0.55),
                 desc, size=13, color=TEXT)
        x_feat += feat_w + feat_gap

    footer(s)


def slide_dashboard_split(prs, *, page, total):
    s = new_slide(prs)
    header(s, "Tablouri de bord personalizate pe rol",
           subtitle="Vedere adaptată automat la rolul utilizatorului autentificat",
           page=page, total=total)
    add_image_safe(s, "10-buyer-dashboard.png",
                   Inches(0.5), Inches(1.85), Inches(6.0), Inches(4.15),
                   caption="Tablou de bord CUMPĂRĂTOR — oferte, salvări, portofel")
    add_image_safe(s, "20-seller-dashboard.png",
                   Inches(6.8), Inches(1.85), Inches(6.0), Inches(4.15),
                   caption="Tablou de bord VÂNZĂTOR — anunțuri, vânzări, câștiguri")
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55),
             "Comparator dedicat pentru 2-5 vehicule — evidențiază cea mai bună valoare în fiecare categorie.",
             size=16, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)


def slide_thanks(prs, *, page, total):
    s = new_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, Inches(5.67), Inches(2.3), Inches(2), Inches(0.14), fill=ACCENT)
    add_text(s, 0, Inches(2.65), SLIDE_W, Inches(1.5),
             "Vă mulțumesc!", size=80, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.2), SLIDE_W, Inches(0.9),
             "Întrebări?", size=40, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(5.45), SLIDE_W, Inches(0.6),
             "Cioara Raul  •  Licență 2026",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(6.05), SLIDE_W, Inches(0.5),
             "Universitatea din Oradea  •  Specializarea Calculatoare",
             size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)


# -------------- asamblare prezentare --------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 18 slide-uri total (1 = copertă, 18 = mulțumiri)
    TOTAL = 18

    # 1. Copertă
    slide_title(prs)

    # 2. Stiva tehnologică
    slide_tech_grid(prs, page=2, total=TOTAL)

    # 3. Securitate și servicii externe
    slide_text(prs, "Securitate și servicii externe",
        [
            "Jetoane JWT cu 2 secrete: acces (15 min) + reîmprospătare (7 zile).",
            "Listă neagră de jetoane pentru invalidarea sesiunilor la deconectare.",
            "Parole cifrate cu scrypt și salt aleatoriu.",
            "Limitare: max 5 încercări de autentificare pe minut.",
            "OAuth2 prin Passport.js (Google și Facebook).",
            "Cloudinary pentru imagini, Mailtrap SMTP pentru trimitere email.",
        ], page=3, total=TOTAL,
        subtitle="Cum am asigurat protecția datelor utilizatorilor")

    # 4. Arhitectura aplicației
    slide_architecture(prs, page=4, total=TOTAL)

    # 5. Modelul de date
    slide_erd(prs, page=5, total=TOTAL)

    # 6. Module backend
    slide_modules(prs, page=6, total=TOTAL)

    # 7. Tipuri de utilizatori
    slide_four_roles(prs, page=7, total=TOTAL)

    # 8. Cerințe funcționale și nefuncționale
    slide_two_columns(prs, "Cerințe funcționale și tehnice",
        "Funcționale",
        [
            "Autentificare locală și prin furnizori externi.",
            "Operații complete pe anunțuri.",
            "Sistem de oferte cu stări multiple.",
            "Mesagerie cu urmărire mesaje necitite.",
            "Portofel cu sume blocate.",
            "Recomandări vehicule similare.",
        ],
        "Tehnice",
        [
            "Securitate: JWT, scrypt, CORS.",
            "Limitare pe rutele sensibile.",
            "Performanță: paginare și randare pe server.",
            "Disponibilitate: SQLite în mod WAL.",
            "Uzabilitate: responsiv, interfață în română.",
            "Mentenabilitate: arhitectură modulară.",
        ], page=8, total=TOTAL)

    # 9. Autentificare
    slide_auth(prs, page=9, total=TOTAL)

    # 10. Listare anunțuri
    slide_big_image(prs, "Listarea și filtrarea anunțurilor",
        "11-find-vehicle.png", page=10, total=TOTAL,
        subtitle="Filtrare după marcă, model, preț, an, combustibil, transmisie și localitate",
        caption="200 de vehicule disponibile din 10 producători diferiți")

    # 11. Detaliu vehicul
    slide_big_image(prs, "Pagina de detaliu vehicul",
        "12-vehicle-detail.png", page=11, total=TOTAL,
        subtitle="Peste 40 de câmpuri tehnice, galerie de imagini și contact direct cu vânzătorul",
        caption="Pagina de detaliu cu informații complete și vehicule similare sugerate")

    # 12. Sistemul de oferte
    slide_with_image(prs, "Sistemul de oferte",
        "15-bids.png",
        [
            "Ofertă cu sumă și mesaj de negociere.",
            "Stări: în așteptare, acceptată, respinsă, retrasă.",
            "La acceptare → vehiculul devine VÂNDUT.",
            "Suma cumpărătorului este blocată în portofel.",
            "Doar rolul CUMPĂRĂTOR poate plasa oferte.",
        ], page=12, total=TOTAL,
        image_caption="Oferte active și istoric complet")

    # 13. Mesagerie
    slide_with_image(prs, "Mesagerie cumpărător — vânzător",
        "16b-messages-thread.png",
        [
            "Conversații directe între cumpărător și vânzător.",
            "Editare și ștergere mesaje individuale.",
            "Ștergere completă a conversației.",
            "Urmărire separată a mesajelor necitite per participant.",
            "Inițiere din pagina de detaliu a vehiculului.",
        ], page=13, total=TOTAL,
        image_caption="Fir complet de conversație",
        image_left=True)

    # 14. Portofel
    slide_with_image(prs, "Portofel digital și tranzacții",
        "17-wallet.png",
        [
            "Sold disponibil și sold blocat.",
            "Card asociat cu depozite și retrageri.",
            "Tipuri: depozit, retragere, plată.",
            "Istoric paginat și filtrat pe tip/stare.",
            "Comision de 5% reținut automat la vânzare.",
        ], page=14, total=TOTAL,
        image_caption="Sold, card și istoric tranzacții")

    # 15. Tablouri de bord
    slide_dashboard_split(prs, page=15, total=TOTAL)

    # 16. Vehicule similare
    slide_text(prs, "Recomandări de vehicule similare",
        [
            "Rută dedicată: GET /vehicule/:id/similare.",
            "Filtrare după aceeași marcă și interval de preț ±15%.",
            "Sortare după distanța absolută față de prețul de referință.",
            "Returnează primele 5 vehicule pe pagina de detaliu.",
            "Extensibil cu algoritm de învățare automată (KNN, vectori).",
        ], page=16, total=TOTAL,
        subtitle="„Vehicule asemănătoare\" sugerate pe pagina de detaliu")

    # 17. Concluzii și îmbunătățiri viitoare
    slide_two_columns(prs, "Concluzii și îmbunătățiri viitoare",
        "Ce am realizat",
        [
            "Platformă web completă și funcțională.",
            "Autentificare OAuth și verificare email.",
            "Sistem complet de oferte.",
            "Mesagerie directă între utilizatori.",
            "Portofel cu tranzacții și sume blocate.",
            "Recomandări de vehicule similare.",
        ],
        "Îmbunătățiri viitoare",
        [
            "Traducere completă a componentelor (i18n).",
            "Integrare plată prin Stripe.",
            "Arhitectură mai curată pe straturi.",
            "Aplicație mobilă React Native.",
            "Implementare completă a panoului de administrare (ADMIN).",
            "Notificări în timp real.",
        ], page=17, total=TOTAL)

    # 18. Mulțumiri
    slide_thanks(prs, page=18, total=TOTAL)

    prs.save(OUT)
    print(f"Salvat: {OUT}")


if __name__ == "__main__":
    build()
